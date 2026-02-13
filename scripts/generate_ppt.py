# Capture all app screens with Playwright and build a PPTX
# Usage: python scripts/generate_ppt.py

import time
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

from playwright.sync_api import sync_playwright

ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = ROOT / "exports"
SCREEN_DIR = EXPORT_DIR / "screenshots"
PPTX_PATH = EXPORT_DIR / "Weather_Dashboard_Screens.pptx"

LOGIN = {"username": "admin", "password": "forecast2026"}
BASE_URL = os.environ.get("APP_URL", "http://127.0.0.1:5000")

SCREENS = [
    {"title": "Login", "action": "goto", "path": "/login", "notes": "Sign-in screen. Use provided demo credentials (admin / forecast2026)."},
    {"title": "Dashboard — Overview", "action": "click", "selector": ".nav-item[data-page='dashboard']", "post": "#overview-subtab", "notes": "Overview KPIs, temperature map and active alerts."},
    {"title": "Dashboard — Trends & Cities", "action": "click_subtab", "sidebar_selector": ".nav-item[data-page='dashboard']", "subtab_selector": "button.dashboard-tab[data-subtab='trends']", "post": "#trends-subtab", "notes": "Temperature trends, cities overview and weekly summary."},
    {"title": "Dashboard — Historical Data", "action": "click_subtab", "sidebar_selector": ".nav-item[data-page='dashboard']", "subtab_selector": "button.dashboard-tab[data-subtab='historical']", "post": "#historical-subtab", "notes": "2-year analysis, heatmap and year comparisons."},
    {"title": "Dashboard — AI Predictions", "action": "click_subtab", "sidebar_selector": ".nav-item[data-page='dashboard']", "subtab_selector": "button.dashboard-tab[data-subtab='predictions']", "post": "#predictions-subtab", "notes": "Demand predictions, energy estimates and seasonal patterns."},
    {"title": "Analytics", "action": "click", "selector": ".nav-item[data-page='analytics']", "post": "#analyticsPage", "notes": "Advanced analytics, heatmaps and distribution charts."},
    {"title": "Forecasting", "action": "click", "selector": ".nav-item[data-page='forecast']", "post": "#forecastPage", "notes": "Forecasting tools and model outputs."},
    {"title": "Cities", "action": "click", "selector": ".nav-item[data-page='cities']", "post": "#citiesPage", "notes": "City-level overview and configuration."},
    {"title": "Alerts", "action": "click", "selector": ".nav-item[data-page='alerts']", "post": "#alertsPage", "notes": "Active alerts panel and acknowledgment controls."},
    {"title": "Insights", "action": "click", "selector": ".nav-item[data-page='insights']", "post": "#insightsPage", "notes": "Actionable insights and recommendations."},
    {"title": "Demand Intel", "action": "click", "selector": ".nav-item[data-page='demand-intel']", "post": "#demandIntelPage", "notes": "Demand intelligence and wave sequencing."},
    {"title": "Reports", "action": "click", "selector": ".nav-item[data-page='reports']", "post": "#reportsPage", "notes": "Exportable reports and PDF/Excel controls."},
]

DISPLAY_WIDTH = 1366
DISPLAY_HEIGHT = 768


def ensure_dirs():
    EXPORT_DIR.mkdir(exist_ok=True)
    SCREEN_DIR.mkdir(parents=True, exist_ok=True)


def add_image_slide(prs, title, image_path, notes_text=""):
    slide_layout = prs.slide_layouts[5]  # Title + content (blank-like)
    slide = prs.slides.add_slide(slide_layout)
    # Title
    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(28)

    # Add image centered, keep aspect ratio
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    slide.shapes.add_picture(str(image_path), left, top, width=width)

    # Speaker notes
    if notes_text:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes_text


def build_pptx(screenshots_with_notes):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "ForecastWell — UI Screenshots"
    slide.placeholders[1].text = "Auto-generated slides: each tab / sub-tab with brief notes."

    for item in screenshots_with_notes:
        add_image_slide(prs, item['title'], item['path'], item.get('notes', ''))

    prs.save(PPTX_PATH)
    print(f"Saved PPTX to: {PPTX_PATH}")
    return PPTX_PATH


def wait_for_server(page, url, timeout=20):
    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            resp = page.goto(url, wait_until="networkidle", timeout=5000)
            if resp and resp.ok:
                return True
        except Exception:
            time.sleep(0.5)
    raise RuntimeError(f"App did not respond at {url} within {timeout}s")


def capture():
    ensure_dirs()
    screenshots = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(viewport={"width": DISPLAY_WIDTH, "height": DISPLAY_HEIGHT})
        page = context.new_page()

        # Wait for server
        print(f"Connecting to {BASE_URL}/login ...")
        wait_for_server(page, f"{BASE_URL}/login", timeout=30)

        # Login
        page.goto(f"{BASE_URL}/login")
        page.fill('#username', LOGIN['username'])
        page.fill('#password', LOGIN['password'])
        page.click('#loginBtn')
        page.wait_for_url(f"{BASE_URL}/", timeout=8000)
        time.sleep(1)

        # Capture per SCREENS
        for i, s in enumerate(SCREENS):
            title = s['title']
            print(f"Capturing: {title}")

            if s['action'] == 'goto':
                page.goto(f"{BASE_URL}{s['path']}", wait_until='networkidle')
            elif s['action'] == 'click':
                page.click(s['selector'])
                # wait for a page-specific element if post is provided
                if s.get('post'):
                    try:
                        page.wait_for_selector(s['post'], timeout=4000)
                    except Exception:
                        time.sleep(0.6)
                else:
                    time.sleep(0.6)
            elif s['action'] == 'click_subtab':
                # Ensure main sidebar clicked
                if s.get('sidebar_selector'):
                    page.click(s['sidebar_selector'])
                    time.sleep(0.3)
                page.click(s['subtab_selector'])
                if s.get('post'):
                    try:
                        page.wait_for_selector(s['post'], timeout=4000)
                    except Exception:
                        time.sleep(0.6)
                else:
                    time.sleep(0.6)

            # Small delay for animations
            time.sleep(0.6)

            fname = SCREEN_DIR / f"{i+1:02d}_{title.replace(' ', '_').replace('—','').replace('&','and')}.png"
            page.screenshot(path=str(fname), full_page=True)
            screenshots.append({"title": title, "path": fname, "notes": s.get('notes', '')})

        browser.close()

    # Build PPTX
    pptx_path = build_pptx(screenshots)
    return pptx_path


if __name__ == '__main__':
    try:
        out = capture()
        print("Finished: ", out)
    except Exception as e:
        print("Error during capture:", e)
        raise
