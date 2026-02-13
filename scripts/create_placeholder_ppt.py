from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = ROOT / "exports"
EXPORT_DIR.mkdir(exist_ok=True)
PPTX_PATH = EXPORT_DIR / "Weather_Dashboard_Placeholders.pptx"

slides = [
    {"title": "ForecastWell — Weather-based Demand Dashboard", "note": "Intro: ForecastWell visualizes weather-driven demand for HVAC & consumer durables. This deck contains UI slides with speaker notes.", "placeholder": "App logo / Home screenshot"},
    {"title": "Login", "note": "Demo credentials: admin / forecast2026. Shows login UX and success overlay.", "placeholder": "01_Login.png"},
    {"title": "Dashboard — Overview", "note": "KPIs: Hottest Day / Night, Temperature Map, Active Alerts, Wave Sequence (NOW → +2w → +6w).", "placeholder": "02_Dashboard_Overview.png"},
    {"title": "Dashboard — Trends & Cities", "note": "Temperature trends selector, cities overview cards, weekly summary for planning.", "placeholder": "03_Dashboard_Trends_Cities.png"},
    {"title": "Dashboard — Historical Data", "note": "Two-year analysis, monthly heatmap and year-over-year comparisons.", "placeholder": "04_Dashboard_Historical.png"},
    {"title": "Dashboard — AI Predictions", "note": "Demand forecasts, energy estimates and seasonal demand patterns to inform inventory.", "placeholder": "05_Dashboard_Predictions.png"},
    {"title": "Analytics", "note": "Advanced analytics: KPI rings, heatmaps, radar and distribution charts for deep insights.", "placeholder": "06_Analytics.png"},
    {"title": "Forecasting", "note": "Scenario-driven forecast controls and model outputs for planning horizons.", "placeholder": "07_Forecasting.png"},
    {"title": "Cities", "note": "Per-city metadata, demand index, and city-level forecasts for prioritization.", "placeholder": "08_Cities.png"},
    {"title": "Alerts", "note": "Active alerts with priorities and recommended actions (DSB methodology).", "placeholder": "09_Alerts.png"},
    {"title": "Insights", "note": "Actionable recommendations (stock, promos) prioritized by impact/urgency.", "placeholder": "10_Insights.png"},
    {"title": "Demand Intel", "note": "Wave sequencing and leading indicators (Wave 1/2/3) to drive operations.", "placeholder": "11_DemandIntel.png"},
    {"title": "Reports", "note": "Export options (PDF/Excel) and scheduled reporting for stakeholders.", "placeholder": "12_Reports.png"},
    {"title": "Top header — Utilities", "note": "Global search, date range, live indicator, and export/upload quick actions.", "placeholder": "13_Header_Controls.png"},
    {"title": "Next steps & Demo checklist", "note": "Demo flow: login → Overview KPIs → Alerts → Predictions → Q&A. Replace placeholders with screenshots before presenting.", "placeholder": "14_Closing.png"},
]

prs = Presentation()
# Title slide
title_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_layout)
slide.shapes.title.text = slides[0]["title"]
slide.placeholders[1].text = "UI walkthrough — tabs, key screens & talking points"
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = slides[0]["note"]

# Content slides
for s in slides[1:]:
    layout = prs.slide_layouts[5]  # Title + content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = s["title"]

    # Add placeholder shape for screenshot
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(4.8)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    shape.fill.background()  # keep transparent
    # Add a caption text inside the placeholder
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"[Add screenshot: {s['placeholder']}]"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.italic = False
    p.alignment = 1  # centered

    # Speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = s["note"]

prs.save(PPTX_PATH)
print(f"Created placeholder PPTX: {PPTX_PATH}")
